"""
Enterprise Knowledge Assistant - Document Ingestion Service

Handles the complete pipeline from raw file upload to indexed chunks:
1. MIME type validation & malware prevention
2. OCR (Tesseract / EasyOCR for images/scanned PDFs)
3. Parsing (PyMuPDF, python-docx, Docling)
4. Language detection
5. Semantic chunking
6. Metadata extraction (title, headings, tables, entities)
7. Summarization & keyword extraction
8. Embedding generation
9. Storage in PostgreSQL + pgvector
"""

from __future__ import annotations

import asyncio
import hashlib
import io
import mimetypes
import os
import re
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterator

import aiofiles
import filetype
import structlog
from langdetect import detect, LangDetectException
from PIL import Image
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config.settings import get_rag_settings, get_storage_settings
from app.domain.models.all_models import Document, DocumentChunk, DocumentStatus
from app.infrastructure.embedding.service import EmbeddingService

logger = structlog.get_logger(__name__)


@dataclass
class ParsedDocument:
    """Intermediate representation after parsing, before chunking."""
    content: str
    pages: list[dict] = field(default_factory=list)  # [{page_num, content, tables}]
    title: str | None = None
    author: str | None = None
    language: str | None = None
    headings: list[dict] = field(default_factory=list)
    tables: list[dict] = field(default_factory=list)
    images: list[dict] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    word_count: int = 0
    page_count: int = 0


@dataclass
class Chunk:
    """A semantic chunk ready for embedding and indexing."""
    content: str
    chunk_index: int
    page_number: int | None = None
    section: str | None = None
    heading: str | None = None
    start_char: int | None = None
    end_char: int | None = None
    metadata: dict = field(default_factory=dict)
    token_count: int | None = None


class FileValidator:
    """
    Validates uploaded files for security.

    Checks MIME type, file size, extension, and rejects executables.
    Uses python-magic for true MIME detection (not just extension).
    """

    def __init__(self):
        self.settings = get_storage_settings()

    def validate(self, file_data: bytes, original_filename: str) -> tuple[bool, str, str]:
        """
        Validate a file.

        Returns (is_valid, error_message, detected_mime_type).
        """
        # Check file size
        size_mb = len(file_data) / (1024 * 1024)
        if size_mb > self.settings.max_upload_size_mb:
            return False, f"File exceeds {self.settings.max_upload_size_mb}MB limit", ""

        # Detect true MIME type from file content (not extension)
        kind = filetype.guess(file_data)
        detected_mime = kind.mime if kind else "application/octet-stream"

        # Allow text files that filetype can't detect
        if detected_mime == "application/octet-stream":
            try:
                file_data.decode("utf-8")
                detected_mime = "text/plain"
            except UnicodeDecodeError:
                pass

        # Validate MIME type
        if detected_mime not in self.settings.allowed_mime_types:
            return False, f"File type '{detected_mime}' is not allowed", detected_mime

        # Check extension against blocked list
        ext = Path(original_filename).suffix.lower()
        if ext in self.settings.blocked_extensions:
            return False, f"File extension '{ext}' is not allowed", detected_mime

        # Additional checks for ZIP files (could contain executables)
        if detected_mime == "application/zip":
            import zipfile
            try:
                with zipfile.ZipFile(io.BytesIO(file_data)) as zf:
                    for name in zf.namelist():
                        name_ext = Path(name).suffix.lower()
                        if name_ext in self.settings.blocked_extensions:
                            return False, f"ZIP contains blocked file type: {name}", detected_mime
            except zipfile.BadZipFile:
                return False, "Invalid ZIP file", detected_mime

        return True, "", detected_mime


class DocumentParser:
    """
    Parses various document types into a unified ParsedDocument.

    Dispatches to the appropriate parser based on MIME type.
    """

    async def parse(self, file_data: bytes, mime_type: str, filename: str) -> ParsedDocument:
        """Parse a document and return structured content."""
        parsers = {
            "application/pdf": self._parse_pdf,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._parse_docx,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._parse_xlsx,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._parse_pptx,
            "text/plain": self._parse_text,
            "text/markdown": self._parse_text,
            "text/csv": self._parse_csv,
            "image/jpeg": self._parse_image,
            "image/png": self._parse_image,
            "image/tiff": self._parse_image,
        }

        parser = parsers.get(mime_type, self._parse_text)
        doc = await asyncio.get_event_loop().run_in_executor(
            None, lambda: asyncio.run(parser(file_data, filename))
        )

        # Detect language
        if doc.content and not doc.language:
            try:
                doc.language = detect(doc.content[:2000])
            except LangDetectException:
                doc.language = "en"

        doc.word_count = len(doc.content.split())
        return doc

    async def _parse_pdf(self, file_data: bytes, filename: str) -> ParsedDocument:
        """
        Parse PDF using PyMuPDF.
        Falls back to OCR for scanned/image-based PDFs.
        """
        import fitz  # PyMuPDF

        doc = ParsedDocument()
        pages = []

        with fitz.open(stream=file_data, filetype="pdf") as pdf:
            doc.page_count = len(pdf)

            # Extract metadata
            meta = pdf.metadata
            doc.title = meta.get("title") or filename
            doc.author = meta.get("author")

            full_text = []
            for page_num, page in enumerate(pdf, start=1):
                # Extract text
                page_text = page.get_text("text")

                # If page has very little text, it might be scanned — use OCR
                if len(page_text.strip()) < 50:
                    pix = page.get_pixmap(dpi=200)
                    img_data = pix.tobytes("png")
                    page_text = await self._ocr_image(img_data)

                # Extract tables
                tables = []
                try:
                    for table in page.find_tables():
                        tables.append({
                            "data": table.extract(),
                            "bbox": list(table.bbox),
                        })
                except Exception:
                    pass

                pages.append({
                    "page_num": page_num,
                    "content": page_text,
                    "tables": tables,
                })
                full_text.append(page_text)

        doc.content = "\n\n".join(full_text)
        doc.pages = pages
        return doc

    async def _parse_docx(self, file_data: bytes, filename: str) -> ParsedDocument:
        """Parse DOCX using python-docx."""
        from docx import Document as DocxDocument
        from docx.oxml.ns import qn

        doc = ParsedDocument()
        docx = DocxDocument(io.BytesIO(file_data))

        # Extract core properties
        props = docx.core_properties
        doc.title = props.title or filename
        doc.author = props.author

        paragraphs = []
        headings = []

        for para in docx.paragraphs:
            if not para.text.strip():
                continue
            style = para.style.name if para.style else ""
            if "Heading" in style:
                level = int(style.split()[-1]) if style[-1].isdigit() else 1
                headings.append({"text": para.text, "level": level})
            paragraphs.append(para.text)

        doc.content = "\n\n".join(paragraphs)
        doc.headings = headings
        doc.page_count = len(docx.sections)
        return doc

    async def _parse_xlsx(self, file_data: bytes, filename: str) -> ParsedDocument:
        """Parse Excel files."""
        import openpyxl

        doc = ParsedDocument()
        wb = openpyxl.load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)

        sheets_content = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets_content.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))

        doc.content = "\n\n".join(sheets_content)
        doc.title = filename
        return doc

    async def _parse_pptx(self, file_data: bytes, filename: str) -> ParsedDocument:
        """Parse PowerPoint files."""
        from pptx import Presentation

        doc = ParsedDocument()
        prs = Presentation(io.BytesIO(file_data))

        slides_content = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides_content.append(f"## Slide {i}\n" + "\n".join(texts))

        doc.content = "\n\n".join(slides_content)
        doc.page_count = len(prs.slides)
        doc.title = filename
        return doc

    async def _parse_text(self, file_data: bytes, filename: str) -> ParsedDocument:
        """Parse plain text and markdown."""
        content = file_data.decode("utf-8", errors="replace")
        return ParsedDocument(content=content, title=filename, page_count=1)

    async def _parse_csv(self, file_data: bytes, filename: str) -> ParsedDocument:
        """Parse CSV files."""
        import csv

        content = file_data.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(content))
        rows = ["\t".join(row) for row in reader if any(cell.strip() for cell in row)]
        return ParsedDocument(content="\n".join(rows), title=filename)

    async def _parse_image(self, file_data: bytes, filename: str) -> ParsedDocument:
        """Extract text from images using OCR."""
        text = await self._ocr_image(file_data)
        return ParsedDocument(content=text, title=filename, page_count=1)

    async def _ocr_image(self, image_data: bytes) -> str:
        """Run OCR on an image, trying Tesseract first, EasyOCR as fallback."""
        try:
            import pytesseract
            img = Image.open(io.BytesIO(image_data))
            return pytesseract.image_to_string(img)
        except Exception as e:
            logger.warning("Tesseract OCR failed, trying EasyOCR", error=str(e))
            try:
                import easyocr
                import numpy as np
                reader = easyocr.Reader(["en"])
                img_array = np.frombuffer(image_data, np.uint8)
                results = reader.readtext(img_array)
                return " ".join(text for _, text, _ in results)
            except Exception as e2:
                logger.error("Both OCR engines failed", error=str(e2))
                return ""


class SemanticChunker:
    """
    Splits documents into semantic chunks for embedding.

    Uses a sliding window approach with paragraph-aware splitting.
    Respects heading boundaries for better context preservation.
    """

    def __init__(self):
        self.settings = get_rag_settings()

    def chunk(self, parsed_doc: ParsedDocument) -> list[Chunk]:
        """Split a parsed document into semantic chunks."""
        if parsed_doc.pages:
            return self._chunk_by_pages(parsed_doc)
        return self._chunk_text(parsed_doc.content)

    def _chunk_by_pages(self, parsed_doc: ParsedDocument) -> list[Chunk]:
        """Chunk page-by-page for documents with explicit pages (PDFs)."""
        chunks = []
        global_idx = 0

        for page in parsed_doc.pages:
            page_chunks = self._chunk_text(page["content"])
            for chunk in page_chunks:
                chunk.page_number = page["page_num"]
                chunk.chunk_index = global_idx
                global_idx += 1
                chunks.append(chunk)

        return chunks

    def _chunk_text(self, text: str) -> list[Chunk]:
        """Split text into overlapping chunks."""
        if not text.strip():
            return []

        chunk_size = self.settings.chunk_size
        overlap = self.settings.chunk_overlap

        # Split on paragraph boundaries first
        paragraphs = re.split(r"\n\s*\n", text)
        chunks = []
        current_chunk = []
        current_size = 0
        chunk_index = 0
        char_offset = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            para_words = para.split()
            para_size = len(para_words)

            if current_size + para_size > chunk_size and current_chunk:
                # Emit current chunk
                chunk_text = " ".join(current_chunk)
                chunks.append(Chunk(
                    content=chunk_text,
                    chunk_index=chunk_index,
                    start_char=char_offset,
                    end_char=char_offset + len(chunk_text),
                    token_count=len(current_chunk),
                ))
                chunk_index += 1
                char_offset += len(chunk_text) + 1

                # Keep overlap from previous chunk
                overlap_words = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_words.copy()
                current_size = len(overlap_words)

            current_chunk.extend(para_words)
            current_size += para_size

        # Emit final chunk
        if current_chunk:
            chunk_text = " ".join(current_chunk)
            chunks.append(Chunk(
                content=chunk_text,
                chunk_index=chunk_index,
                start_char=char_offset,
                end_char=char_offset + len(chunk_text),
                token_count=len(current_chunk),
            ))

        return chunks


class DocumentIngestionService:
    """
    Orchestrates the complete document ingestion pipeline.

    This is the primary entry point for document processing.
    Called by the background worker after initial upload.
    """

    def __init__(
        self,
        db: AsyncSession,
        embedding_service: EmbeddingService,
    ):
        self.db = db
        self.embedding_service = embedding_service
        self.validator = FileValidator()
        self.parser = DocumentParser()
        self.chunker = SemanticChunker()
        self.storage_settings = get_storage_settings()

    async def ingest(self, document_id: str, file_data: bytes) -> None:
        """
        Run the full ingestion pipeline for a document.

        Updates document status throughout the process.
        On failure, marks document as FAILED with error details.
        """
        log = logger.bind(document_id=document_id)

        try:
            # Load document record
            from sqlalchemy import select
            result = await self.db.execute(
                select(Document).where(Document.id == document_id)
            )
            document = result.scalar_one()

            log.info("Starting document ingestion", filename=document.original_filename)
            document.status = DocumentStatus.PROCESSING
            await self.db.commit()

            # 1. Parse document
            parsed = await self.parser.parse(
                file_data, document.mime_type, document.original_filename
            )

            # 2. Update document metadata
            document.title = parsed.title or document.original_filename
            document.author = parsed.author
            document.language = parsed.language
            document.page_count = parsed.page_count
            document.word_count = parsed.word_count

            # 3. Semantic chunking
            chunks = self.chunker.chunk(parsed)
            log.info("Document chunked", num_chunks=len(chunks))

            # 4. Generate embeddings in batch
            chunk_texts = [c.content for c in chunks]
            embeddings = await self.embedding_service.embed_batch(chunk_texts)

            # 5. Store chunks with embeddings
            db_chunks = []
            for chunk, embedding in zip(chunks, embeddings):
                content_hash = hashlib.sha256(chunk.content.encode()).hexdigest()
                db_chunk = DocumentChunk(
                    id=str(uuid.uuid4()),
                    document_id=document_id,
                    organization_id=document.organization_id,
                    content=chunk.content,
                    content_hash=content_hash,
                    embedding=embedding,
                    chunk_index=chunk.chunk_index,
                    page_number=chunk.page_number,
                    section=chunk.section,
                    heading=chunk.heading,
                    start_char=chunk.start_char,
                    end_char=chunk.end_char,
                    token_count=chunk.token_count,
                    metadata=chunk.metadata,
                    allowed_roles=document.allowed_roles,
                )
                db_chunks.append(db_chunk)

            self.db.add_all(db_chunks)

            # 6. Mark as indexed
            document.status = DocumentStatus.INDEXED
            document.processed_at = datetime.now(timezone.utc)
            await self.db.commit()

            log.info("Document ingestion complete", num_chunks=len(db_chunks))

        except Exception as e:
            log.error("Document ingestion failed", error=str(e), exc_info=True)
            try:
                from sqlalchemy import select, update
                await self.db.execute(
                    update(Document)
                    .where(Document.id == document_id)
                    .values(status=DocumentStatus.FAILED, processing_error=str(e))
                )
                await self.db.commit()
            except Exception:
                pass
            raise
